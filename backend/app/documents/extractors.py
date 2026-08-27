"""
Content extractor implementations for the document ingestion pipeline.

- TxtContentExtractor: Reads plain text files directly.
- PdfContentExtractor: Placeholder — raises NotImplementedError (Phase 1).
- DocxContentExtractor: Placeholder — raises NotImplementedError (Phase 1).

The architecture accommodates future image formats (PNG, JPEG) by adding
new extractor classes without modifying existing ones or the pipeline flow.
"""

from pathlib import Path

from app.errors.document_errors import ContentExtractionError


class TxtContentExtractor:
    """Extracts content from plain-text (.txt) files.

    Satisfies the ContentExtractor protocol via structural subtyping.
    """

    async def extract(self, file_path: str, file_type: str) -> str:
        """Read the file at *file_path* and return its text content.

        Args:
            file_path: Absolute or relative path to the .txt file.
            file_type: MIME type or extension identifier (expected: "txt" or "text/plain").

        Returns:
            The full text content of the file.

        Raises:
            ContentExtractionError: If reading the file fails.
        """
        path = Path(file_path)
        if not path.exists():
            raise ContentExtractionError(
                file_name=path.name,
                message=f"File not found: {file_path}",
            )
        try:
            return path.read_text(encoding="utf-8")
        except (OSError, UnicodeDecodeError) as exc:
            raise ContentExtractionError(
                file_name=path.name,
                message=f"Failed to read text file: {path.name}",
                detail=str(exc),
            ) from exc


class PdfContentExtractor:
    """Extracts content from PDF files using PyMuPDF (fitz)."""

    async def extract(self, file_path: str, file_type: str) -> str:
        """Extract text from a .pdf file, page by page.

        Args:
            file_path: Path to the .pdf file.
            file_type: MIME type or extension identifier.

        Returns:
            The full text content with pages separated by page break delimiters.

        Raises:
            ContentExtractionError: If reading the file fails.
        """
        path = Path(file_path)
        if not path.exists():
            raise ContentExtractionError(
                file_name=path.name,
                message=f"File not found: {file_path}",
            )
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            page_texts: list[str] = []
            for page in doc:
                text = page.get_text().strip()
                if text:
                    page_texts.append(text)
            doc.close()

            if not page_texts:
                raise ContentExtractionError(
                    file_name=path.name,
                    message="PDF contains no extractable text (may be image-only).",
                )

            return "\n---PAGE_BREAK---\n".join(page_texts)

        except ContentExtractionError:
            raise
        except Exception as exc:
            raise ContentExtractionError(
                file_name=path.name,
                message=f"Failed to extract PDF content: {path.name}",
                detail=str(exc),
            ) from exc


class DocxContentExtractor:
    """Extracts content from DOCX files using python-docx."""

    async def extract(self, file_path: str, file_type: str) -> str:
        """Extract text from a .docx file.

        Args:
            file_path: Path to the .docx file.
            file_type: MIME type or extension identifier.

        Returns:
            The full text content of the document.

        Raises:
            ContentExtractionError: If reading the file fails.
        """
        path = Path(file_path)
        if not path.exists():
            raise ContentExtractionError(
                file_name=path.name,
                message=f"File not found: {file_path}",
            )
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except Exception as exc:
            raise ContentExtractionError(
                file_name=path.name,
                message=f"Failed to extract DOCX content: {path.name}",
                detail=str(exc),
            ) from exc


class PptxContentExtractor:
    """Extracts content from PowerPoint (.pptx) files using python-pptx.

    Extracts text from:
    - Slide text boxes and shapes
    - Table cells
    - Speaker notes

    Each slide is separated by a page break delimiter for chunk boundary detection.
    """

    async def extract(self, file_path: str, file_type: str) -> str:
        """Extract text from a .pptx file, slide by slide.

        Args:
            file_path: Path to the .pptx file.
            file_type: MIME type or extension identifier.

        Returns:
            The full text content with slides separated by page break delimiters.

        Raises:
            ContentExtractionError: If reading the file fails.
        """
        path = Path(file_path)
        if not path.exists():
            raise ContentExtractionError(
                file_name=path.name,
                message=f"File not found: {file_path}",
            )
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slide_texts: list[str] = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                parts: list[str] = []
                parts.append(f"# Slide {slide_num}")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                parts.append(text)

                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                parts.append(row_text)

                # Extract speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Speaker Notes: {notes}]")

                slide_text = "\n".join(parts)
                if slide_text.strip():
                    slide_texts.append(slide_text)

            if not slide_texts:
                raise ContentExtractionError(
                    file_name=path.name,
                    message="PowerPoint file contains no extractable text.",
                )

            # Join slides with page break delimiter (same as PDF processor)
            return "\n---PAGE_BREAK---\n".join(slide_texts)

        except ContentExtractionError:
            raise
        except Exception as exc:
            raise ContentExtractionError(
                file_name=path.name,
                message=f"Failed to extract PowerPoint content: {path.name}",
                detail=str(exc),
            ) from exc
